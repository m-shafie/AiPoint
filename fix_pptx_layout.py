#!/usr/bin/env python3
"""
Repair AI-export PowerPoint issues: 
- Slide size / EMU scaling
- Persian cs font and language tags
- Parentheses size fix (Split-run strategy)
- Global font size override
- Latin (and optional ea) typeface for English text
- Optional paragraph margin stripping and single line spacing
- Robust watermark removal via python-pptx
"""

from __future__ import annotations

import argparse
import io
import re
import sys
import zipfile
import xml.etree.ElementTree as ET
from dataclasses import dataclass
from typing import Iterable

# --- NAMESPACES ---
DRAWINGML = "http://schemas.openxmlformats.org/drawingml/2006/main"
PRESENTATIONML = "http://schemas.openxmlformats.org/presentationml/2006/main"

A = f"{{{DRAWINGML}}}"
P = f"{{{PRESENTATIONML}}}"

XFRM = f"{A}xfrm"
OFF = f"{A}off"
EXT = f"{A}ext"
CHOFF = f"{A}chOff"
CHEXT = f"{A}chExt"
RUN = f"{A}r"
T = f"{A}t"
RPR = f"{A}rPr"
PP = f"{A}p"
PPR = f"{A}pPr"
LNSPC = f"{A}lnSpc"
SPCPCT = f"{A}spcPct"
DEF_PPR = f"{A}defPPr"
LATIN = f"{A}latin"
EA = f"{A}ea"
CS = f"{A}cs"
END_PARA_RPR = f"{A}endParaRPr"
SldSz = f"{P}sldSz"

_PERSIAN_RE = re.compile(
    r"[\u0600-\u06FF\u0750-\u077F\u08A0-\u08FF\uFB50-\uFDFF\uFE70-\uFEFF]"
)

# --- SCALING LOGIC ---

def _scale_int(value: str, factor: float) -> str:
    return str(int(round(int(value) * factor)))

def _scale_off_or_choff(elem: ET.Element, sx: float, sy: float) -> None:
    if "x" in elem.attrib: elem.attrib["x"] = _scale_int(elem.attrib["x"], sx)
    if "y" in elem.attrib: elem.attrib["y"] = _scale_int(elem.attrib["y"], sy)

def _scale_ext_or_chext(elem: ET.Element, sx: float, sy: float) -> None:
    if "cx" in elem.attrib: elem.attrib["cx"] = _scale_int(elem.attrib["cx"], sx)
    if "cy" in elem.attrib: elem.attrib["cy"] = _scale_int(elem.attrib["cy"], sy)

def scale_all_xfrm(root: ET.Element, sx: float, sy: float) -> None:
    for xfrm in root.iter(XFRM):
        for child in xfrm:
            if child.tag == OFF or child.tag == CHOFF:
                _scale_off_or_choff(child, sx, sy)
            elif child.tag == EXT or child.tag == CHEXT:
                _scale_ext_or_chext(child, sx, sy)

# --- TYPOGRAPHY & BRACKET FIX ---

def run_text(run: ET.Element) -> str:
    parts: list[str] = []
    for t in run.iter(T):
        if t.text: parts.append(t.text)
        if t.tail: parts.append(t.tail)
    return "".join(parts)

def is_persian_script(s: str) -> bool:
    return bool(_PERSIAN_RE.search(s))

def _ensure_font_slot(elem: ET.Element, slot_tag: str, font: str) -> None:
    if not font: return
    node = elem.find(slot_tag)
    if node is None:
        node = ET.Element(slot_tag)
        elem.append(node)
    node.set("typeface", font)

def ensure_latin_typeface(elem: ET.Element, font: str) -> None: _ensure_font_slot(elem, LATIN, font)
def ensure_ea_typeface(elem: ET.Element, font: str) -> None: _ensure_font_slot(elem, EA, font)
def ensure_cs_typeface(elem: ET.Element, font: str) -> None: _ensure_font_slot(elem, CS, font)

def ensure_rpr(run: ET.Element) -> ET.Element:
    rpr = run.find(RPR)
    if rpr is None:
        idx = 0
        for i, ch in enumerate(run):
            if ch.tag == T:
                idx = i
                break
        rpr = ET.Element(RPR)
        run.insert(idx, rpr)
    return rpr

def split_runs_by_brackets(para: ET.Element) -> None:
    """Splits runs containing '(' or ')' into multiple runs to fix size logic."""
    pattern = re.compile(r'([()])')
    for run in list(para.findall(RUN)):
        t_node = run.find(T)
        if t_node is None or not t_node.text: continue
        parts = pattern.split(t_node.text)
        if len(parts) > 1:
            idx = list(para).index(run)
            for part in parts:
                if not part: continue
                new_run = ET.fromstring(ET.tostring(run))
                new_run.find(T).text = part
                para.insert(idx, new_run)
                idx += 1
            para.remove(run)

def fix_paragraph_typography(
    para: ET.Element,
    persian_font: str,
    latin_font: str,
    latin_lang: str,
    persian_lang: str,
    cs_on_all_runs: bool,
    set_ea_font: bool,
) -> None:
    split_runs_by_brackets(para)
    runs = [r for r in para.iter(RUN) if r.tag == RUN]
    full = "".join(run_text(r) for r in runs)
    para_has_persian = is_persian_script(full)

    for run in runs:
        text = run_text(run)
        rpr = ensure_rpr(run)
        is_bracket = text in ("(", ")")
        persian = is_persian_script(text)

        if is_bracket or persian:
            rpr.set("lang", persian_lang)
            rpr.set("cs", "1")
            ensure_cs_typeface(rpr, persian_font)
            if latin_font: ensure_latin_typeface(rpr, latin_font)
        else:
            rpr.set("lang", latin_lang)
            rpr.attrib.pop("cs", None)
            if latin_font:
                ensure_latin_typeface(rpr, latin_font)
                if set_ea_font: ensure_ea_typeface(rpr, latin_font)
            if cs_on_all_runs: ensure_cs_typeface(rpr, persian_font)

    for ep in para.iter(END_PARA_RPR):
        if para_has_persian:
            ep.set("lang", persian_lang)
            ep.set("cs", "1")
            ensure_cs_typeface(ep, persian_font)
        else:
            ep.set("lang", latin_lang)
        if latin_font: ensure_latin_typeface(ep, latin_font)

# --- FORMATTING & SIZE ---

def apply_font_size_to_root(root: ET.Element, size_pt: float) -> None:
    sz_val = str(int(size_pt * 100))
    for rpr in root.iter(RPR): rpr.set("sz", sz_val)
    for ep in root.iter(END_PARA_RPR): ep.set("sz", sz_val)
    for def_rpr in root.iter(f"{A}defRPr"): def_rpr.set("sz", sz_val)

def strip_paragraph_margins(root: ET.Element) -> None:
    for ppr in root.iter(PPR):
        for attr in ("marL", "marR", "marT", "marB"): ppr.attrib.pop(attr, None)

def apply_line_spacing_multiple(root: ET.Element, multiple: float) -> None:
    """Set paragraph line spacing as a multiple of single (DrawingML spcPct: 1.0 -> 100000)."""
    val = str(max(1, int(round(float(multiple) * 100000))))
    for ap in root.iter(PP):
        ppr = ap.find(PPR)
        if ppr is None:
            ppr = ET.Element(PPR)
            ap.insert(0, ppr)
        for ln in list(ppr.findall(LNSPC)):
            ppr.remove(ln)
        ln_spc = ET.SubElement(ppr, LNSPC)
        ET.SubElement(ln_spc, SPCPCT, val=val)

def apply_line_spacing_single(root: ET.Element) -> None:
    apply_line_spacing_multiple(root, 1.0)

# --- XML WRAPPERS ---

def process_drawing_xml(data: bytes, sx: float | None, sy: float | None, **kwargs) -> bytes:
    root = ET.fromstring(data)
    if sx is not None and sy is not None: scale_all_xfrm(root, sx, sy)
    if kwargs.get('typography'):
        for para in root.iter(PP):
            fix_paragraph_typography(
                para, kwargs['persian_font'], kwargs['latin_font'],
                kwargs['latin_lang'], kwargs['persian_lang'],
                kwargs['cs_on_all_runs'], kwargs['set_ea_font']
            )
    if kwargs.get('strip_margins'): strip_paragraph_margins(root)
    if kwargs.get('line_spacing_single'):
        apply_line_spacing_multiple(root, 1.0)
    elif kwargs.get('line_spacing_multiple') is not None:
        apply_line_spacing_multiple(root, float(kwargs['line_spacing_multiple']))
    if kwargs.get('font_size'): apply_font_size_to_root(root, kwargs['font_size'])
    return ET.tostring(root, encoding="utf-8", xml_declaration=True)

# --- WATERMARK REMOVAL (FULL) ---

@dataclass
class WatermarkRemovalConfig:
    mode: str; scope: str | None; strict_name: str; strict_box: tuple[int, int, int, int]
    strict_tolerance: int; x_frac: float; y_frac: float; name_substring: str | None; include_placeholders: bool

def apply_watermark_removal(package_bytes: bytes, cfg: WatermarkRemovalConfig) -> tuple[bytes, int]:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.shapes.base import BaseShape
        from pptx.shapes.group import GroupShape
        from pptx.shapes.shapetree import GroupShapes
    except ImportError:
        raise SystemExit("Watermark removal requires python-pptx. Install with: pip install python-pptx")

    def _slide_space_origin(shape: BaseShape) -> tuple[float, float]:
        x, y = float(shape.left), float(shape.top)
        parent = shape._parent
        while isinstance(parent, GroupShapes):
            owner = parent._parent
            if isinstance(owner, GroupShape):
                x += float(owner.left); y += float(owner.top); parent = owner._parent
            else: break
        return x, y

    def _matches(shape: BaseShape, sw: int, sh: int) -> bool:
        if cfg.mode == "strict":
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE or (shape.name or "") != cfg.strict_name: return False
            ox, oy = _slide_space_origin(shape); t = cfg.strict_tolerance
            l, top, w, h = cfg.strict_box
            return abs(ox-l)<=t and abs(oy-top)<=t and abs(int(shape.width)-w)<=t and abs(int(shape.height)-h)<=t
        else:
            if cfg.name_substring and cfg.name_substring.lower() in (shape.name or "").lower(): return True
            if not cfg.include_placeholders and shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER: return False
            ox, oy = _slide_space_origin(shape)
            cx, cy = ox + float(shape.width)/2, oy + float(shape.height)/2
            return cx >= sw * cfg.x_frac and cy >= sh * cfg.y_frac

    def _scan(shapes, sw, sh, do_remove):
        n = 0
        for shape in list(shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                n += _scan(shape.shapes, sw, sh, do_remove)
                if _matches(shape, sw, sh):
                    n += 1; 
                    if do_remove: shape._element.getparent().remove(shape._element)
            elif _matches(shape, sw, sh):
                n += 1
                if do_remove: shape._element.getparent().remove(shape._element)
        return n

    prs = Presentation(io.BytesIO(package_bytes))
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    count = 0
    for master in prs.slide_masters:
        count += _scan(master.shapes, sw, sh, True)
        for layout in master.slide_layouts: count += _scan(layout.shapes, sw, sh, True)
    if cfg.scope == "all" or (cfg.mode == "heuristic" and cfg.scope is None):
        for slide in prs.slides: count += _scan(slide.shapes, sw, sh, True)
    
    out = io.BytesIO(); prs.save(out)
    return out.getvalue(), count

# --- ZIP UTILS & CLI ---

def parse_parts_arg(s: str) -> set[str]:
    allowed = {"slides", "layouts", "masters", "notes", "diagrams", "handouts"}
    parts = {x.strip().lower() for x in s.split(",") if x.strip()}
    bad = parts - allowed
    if bad: raise ValueError(f"Unknown parts: {bad}")
    return parts or allowed

def prefixes_for_parts(parts: set[str]) -> list[str]:
    mapping = {
        "slides": ["ppt/slides/slide"], "layouts": ["ppt/slideLayouts/slideLayout"],
        "masters": ["ppt/slideMasters/slideMaster"], "diagrams": ["ppt/diagrams/"],
        "notes": ["ppt/notesSlides/notesSlide", "ppt/notesMasters/notesMaster"],
        "handouts": ["ppt/handoutMasters/handoutMaster"],
    }
    return [p for k in parts for p in mapping[k]]

def transform_pptx(input_path, output_path, **kwargs):
    pkg_source = input_path
    if kwargs.get('watermark_cfg'):
        with open(input_path, "rb") as f:
            patched, n = apply_watermark_removal(f.read(), kwargs['watermark_cfg'])
            pkg_source = io.BytesIO(patched)

    opener = zipfile.ZipFile(pkg_source, "r")
    with opener as zin:
        pres_data = zin.read("ppt/presentation.xml")
        root_p = ET.fromstring(pres_data)
        sld_sz = root_p.find(f".//{SldSz}")
        src_cx, src_cy = int(sld_sz.get("cx")), int(sld_sz.get("cy"))
        sx, sy = kwargs['target_cx']/src_cx, kwargs['target_cy']/src_cy
        
        sld_sz.set("cx", str(kwargs['target_cx'])); sld_sz.set("cy", str(kwargs['target_cy']))
        if kwargs.get('drop_sldsz_type'): sld_sz.attrib.pop("type", None)

        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            prefixes = prefixes_for_parts(kwargs['parts'])
            for info in zin.infolist():
                data = zin.read(info.filename)
                if info.filename == "ppt/presentation.xml":
                    data = ET.tostring(root_p, encoding="utf-8", xml_declaration=True)
                elif info.filename.endswith(".xml") and any(p in info.filename for p in prefixes):
                    data = process_drawing_xml(data, sx, sy, **kwargs)
                zout.writestr(info, data)

    with open(output_path, "wb") as f: f.write(buf.getvalue())

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("-i", "--input", required=True); ap.add_argument("-o", "--output", required=True)
    ap.add_argument("--target-cx", type=int, default=12192000); ap.add_argument("--target-cy", type=int, default=6858000)
    ap.add_argument("--persian-font", default="B Nazanin"); ap.add_argument("--latin-font", default="Times New Roman")
    ap.add_argument("--font-size", type=float, help="Force font size (pt)")
    ap.add_argument("--strip-paragraph-margins", action="store_true")
    ap.add_argument("--line-spacing-single", action="store_true")
    ap.add_argument("--parts", default="slides,layouts,masters,notes,diagrams,handouts")
    ap.add_argument("--remove-watermark", action="store_true")
    ap.add_argument("--watermark-mode", choices=("strict", "heuristic"), default="strict")
    ap.add_argument("--watermark-strict-box", default="8029180,4844491,1071843,256032")
    ap.add_argument("--watermark-strict-name", default="Image 0")
    args = ap.parse_args()

    w_cfg = None
    if args.remove_watermark:
        box = tuple(map(int, args.watermark_strict_box.split(",")))
        w_cfg = WatermarkRemovalConfig(args.watermark_mode, None, args.watermark_strict_name, box, 8192, 0.55, 0.55, None, False)

    transform_pptx(
        args.input, args.output, target_cx=args.target_cx, target_cy=args.target_cy,
        typography=True, persian_font=args.persian_font, latin_font=args.latin_font,
        latin_lang="en-US", persian_lang="fa-IR", cs_on_all_runs=True, set_ea_font=True,
        font_size=args.font_size, strip_margins=args.strip_paragraph_margins,
        line_spacing_single=args.line_spacing_single, drop_sldsz_type=True,
        parts=parse_parts_arg(args.parts), watermark_cfg=w_cfg
    )

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)