#!/usr/bin/env python3
"""
Simple desktop UI for fix_pptx_layout: queue .pptx files, reorder by drag,
optional merge in list order, then run the same repair pipeline as the CLI.
"""

from __future__ import annotations

import os
import re
import sys
import tempfile
import threading
import tkinter as tk
import zipfile
from copy import deepcopy
from tkinter import filedialog, font, messagebox, ttk

from fix_pptx_layout import (
    WatermarkRemovalConfig,
    parse_parts_arg,
    transform_pptx,
)

try:
    from pptx import Presentation
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from pptx.opc.package import Part
    from pptx.parts.image import ImagePart
    from pptx.oxml.ns import qn
except ImportError as e:  # pragma: no cover - startup guard
    raise SystemExit(
        "This GUI requires python-pptx. Install with: pip install python-pptx"
    ) from e

DEFAULT_PERSIAN = "B Nazanin"
DEFAULT_LATIN = "Times New Roman"
DEFAULT_PARTS = "slides,layouts,masters,notes,diagrams,handouts"


def _sorted_font_families() -> list[str]:
    return sorted({f for f in font.families() if f and not f.startswith("@")})


_SLIDE_PART_RE = re.compile(r"^ppt/slides/slide\d+\.xml$", re.IGNORECASE)


def slide_count_pptx(path: str) -> int:
    """Count slide parts without keeping the file locked (safe alongside COM merge)."""
    with zipfile.ZipFile(path, "r") as zf:
        return sum(
            1 for n in zf.namelist() if _SLIDE_PART_RE.match(n.replace("\\", "/"))
        )


def merge_pptx_files_com(paths: list[str], out_path: str) -> None:
    """Merge decks using PowerPoint; keeps pictures, charts, and masters intact."""
    import pythoncom
    import win32com.client

    com_init = False
    ppt = None
    pres = None
    try:
        pythoncom.CoInitialize()
        com_init = True
        abs_out = os.path.abspath(out_path)
        paths_abs = [os.path.abspath(p) for p in paths]
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.DisplayAlerts = 0
        ppt.Visible = 0
        pres = ppt.Presentations.Open(paths_abs[0], WithWindow=False)
        for extra in paths_abs[1:]:
            n = slide_count_pptx(extra)
            if n <= 0:
                continue
            # Index = slide to insert *after* (MS Learn); append == current slide count.
            pres.Slides.InsertFromFile(extra, pres.Slides.Count, 1, n)
        pres.SaveAs(abs_out)
    finally:
        if pres is not None:
            try:
                pres.Close()
            except Exception:
                pass
        if ppt is not None:
            try:
                ppt.Quit()
            except Exception:
                pass
        if com_init:
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass


def _clear_placeholder_shapes(dest_slide) -> None:
    for shp in list(dest_slide.shapes):
        el = shp.element
        el.getparent().remove(el)


def _clone_image_part_for_slide(dst_slide_part, source_image_part: Part) -> str:
    """Copy *source_image_part* blob into the merged package without Pillow (EMF/WMF/SVG-safe)."""
    pkg = dst_slide_part.package
    ext = source_image_part.partname.ext or "png"
    partname = pkg.next_image_partname(ext)
    new_part = ImagePart(
        partname,
        source_image_part.content_type,
        pkg,
        source_image_part.blob,
        None,
    )
    return dst_slide_part.relate_to(new_part, RT.IMAGE)


def _relink_image_relationships(src_part, dst_part, copied_shape_roots: list) -> None:
    """Point r:embed / r:link in copied XML to new image parts on *dst_part* (same package)."""
    embed = qn("r:embed")
    link = qn("r:link")
    ordered_rids: list[str] = []
    seen: set[str] = set()
    for root in copied_shape_roots:
        for el in root.iter():
            for attr in (embed, link):
                rid = el.get(attr)
                if rid and rid not in seen:
                    seen.add(rid)
                    ordered_rids.append(rid)

    rid_map: dict[str, str] = {}
    for rid in ordered_rids:
        if rid in rid_map:
            continue
        try:
            rel = src_part.rels[rid]
        except KeyError:
            continue
        if rel.is_external or rel.reltype != RT.IMAGE:
            continue
        rid_map[rid] = _clone_image_part_for_slide(dst_part, rel.target_part)

    for root in copied_shape_roots:
        for el in root.iter():
            for attr in (embed, link):
                rid = el.get(attr)
                if rid and rid in rid_map:
                    el.set(attr, rid_map[rid])


def merge_pptx_files_python_relink(paths: list[str], out_path: str) -> None:
    """Append slides with shape XML copy + image relationship relink (no PowerPoint required)."""
    merged = Presentation(paths[0])
    blank_idx = min(6, len(merged.slide_layouts) - 1)
    blank_layout = merged.slide_layouts[blank_idx]
    ext_tag = qn("p:extLst")
    for path in paths[1:]:
        prs = Presentation(path)
        for slide in prs.slides:
            dest = merged.slides.add_slide(blank_layout)
            _clear_placeholder_shapes(dest)
            sp_tree = dest.shapes._spTree
            ext = sp_tree.find(ext_tag)
            copied: list = []
            for shp in slide.shapes:
                new_el = deepcopy(shp.element)
                copied.append(new_el)
                if ext is not None:
                    ext.addprevious(new_el)
                else:
                    sp_tree.append(new_el)
            _relink_image_relationships(slide.part, dest.part, copied)
    merged.save(out_path)


def merge_pptx_files(paths: list[str], out_path: str) -> None:
    """Append slides from paths[1:] onto paths[0]. Prefer PowerPoint COM on Windows."""
    if len(paths) < 2:
        raise ValueError("merge_pptx_files expects at least two paths")
    if sys.platform == "win32":
        try:
            merge_pptx_files_com(paths, out_path)
            return
        except ImportError:
            pass
        except Exception:
            # PowerPoint missing, disabled, or COM error — fall back to python merge.
            pass
    merge_pptx_files_python_relink(paths, out_path)


class PptxLayoutApp(ttk.Frame):
    def __init__(self, master: tk.Tk) -> None:
        super().__init__(master, padding=10)
        self.master.title("PowerPoint Layout Fixer")
        self.master.minsize(720, 420)
        self.files: list[str] = []
        self._drag_index: int | None = None
        self._busy = False

        self.grid(row=0, column=0, sticky="nsew")
        master.rowconfigure(0, weight=1)
        master.columnconfigure(0, weight=1)

        top = ttk.Frame(self)
        top.grid(row=0, column=0, columnspan=2, sticky="ew", pady=(0, 8))
        top.columnconfigure(0, weight=1)

        self.path_var = tk.StringVar(value="")
        ttk.Entry(top, textvariable=self.path_var, state="readonly").grid(
            row=0, column=0, sticky="ew", padx=(0, 6)
        )
        ttk.Button(top, text="Browse", command=self._browse_add).grid(row=0, column=1)

        body = ttk.Panedwindow(self, orient=tk.HORIZONTAL)
        body.grid(row=1, column=0, columnspan=2, sticky="nsew")
        self.rowconfigure(1, weight=1)
        self.columnconfigure(0, weight=1)

        left = ttk.Frame(body, padding=(0, 0, 8, 0))
        right = ttk.Frame(body, padding=(8, 0, 0, 0))
        body.add(left, weight=2)
        body.add(right, weight=1)

        ttk.Label(left, text="Files (top = first in deck merge order)").grid(
            row=0, column=0, columnspan=2, sticky="w"
        )
        ttk.Label(
            left,
            text="Multi-file merge: uses PowerPoint when available (keeps pictures). "
            "Otherwise images are re-linked when possible; complex charts may still need PowerPoint.",
            wraplength=320,
            foreground="#444",
        ).grid(row=1, column=0, columnspan=2, sticky="w", pady=(0, 4))
        self.listbox = tk.Listbox(left, height=14, selectmode=tk.BROWSE, activestyle="dotbox")
        self.listbox.grid(row=2, column=0, sticky="nsew", pady=(4, 4))
        sb = ttk.Scrollbar(left, orient=tk.VERTICAL, command=self.listbox.yview)
        sb.grid(row=2, column=1, sticky="ns", pady=(4, 4))
        self.listbox.config(yscrollcommand=sb.set)
        left.rowconfigure(2, weight=1)
        left.columnconfigure(0, weight=1)

        self.listbox.bind("<ButtonPress-1>", self._lb_press)
        self.listbox.bind("<B1-Motion>", self._lb_motion)
        self.listbox.bind("<ButtonRelease-1>", self._lb_release)

        btns = ttk.Frame(left)
        btns.grid(row=3, column=0, columnspan=2, sticky="w")
        ttk.Button(btns, text="Remove selected", command=self._remove_selected).pack(
            side=tk.LEFT, padx=(0, 6)
        )
        ttk.Button(btns, text="Clear all", command=self._clear_all).pack(side=tk.LEFT)

        ttk.Label(right, text="Options").grid(row=0, column=0, sticky="w")

        opts = ttk.LabelFrame(right, text="Typography & layout", padding=8)
        opts.grid(row=1, column=0, sticky="nsew", pady=(4, 0))
        right.rowconfigure(1, weight=1)
        right.columnconfigure(0, weight=1)

        families = _sorted_font_families()
        r = 0
        ttk.Label(opts, text="Persian font").grid(row=r, column=0, sticky="w")
        self.cb_persian = ttk.Combobox(opts, values=families, width=32, state="normal")
        self.cb_persian.grid(row=r, column=1, sticky="ew", padx=(8, 0), pady=2)
        self.cb_persian.set(DEFAULT_PERSIAN)
        r += 1

        ttk.Label(opts, text="Latin font").grid(row=r, column=0, sticky="w")
        self.cb_latin = ttk.Combobox(opts, values=families, width=32, state="normal")
        self.cb_latin.grid(row=r, column=1, sticky="ew", padx=(8, 0), pady=2)
        self.cb_latin.set(DEFAULT_LATIN)
        r += 1

        ttk.Label(opts, text="Font size (pt, optional)").grid(row=r, column=0, sticky="w")
        self.ent_font_size = ttk.Entry(opts, width=10)
        self.ent_font_size.grid(row=r, column=1, sticky="w", padx=(8, 0), pady=2)
        r += 1

        self.var_watermark = tk.BooleanVar(value=False)
        ttk.Checkbutton(
            opts, text="Remove watermark (strict preset from script)", variable=self.var_watermark
        ).grid(row=r, column=0, columnspan=2, sticky="w", pady=(6, 0))
        r += 1

        ttk.Label(opts, text="Line spacing multiple (optional, e.g. 1, 1.15, 1.5)").grid(
            row=r, column=0, sticky="w", pady=(6, 0)
        )
        self.ent_line_spc = ttk.Entry(opts, width=10)
        self.ent_line_spc.grid(row=r, column=1, sticky="w", padx=(8, 0), pady=(6, 0))
        r += 1
        ttk.Label(
            opts,
            text="Leave font size / line spacing empty to keep original values in the file.",
            wraplength=280,
        ).grid(row=r, column=0, columnspan=2, sticky="w", pady=(8, 0))

        opts.columnconfigure(1, weight=1)

        bottom = ttk.Frame(self)
        bottom.grid(row=2, column=0, columnspan=2, sticky="ew", pady=(10, 0))
        bottom.columnconfigure(0, weight=1)
        self.btn_export = ttk.Button(bottom, text="Export fixed PowerPoint…", command=self._export)
        self.btn_export.pack(side=tk.RIGHT)

    def _browse_add(self) -> None:
        paths = filedialog.askopenfilenames(
            title="Add PowerPoint files",
            filetypes=[("PowerPoint", "*.pptx"), ("All files", "*.*")],
        )
        if not paths:
            return
        for p in paths:
            p = os.path.normpath(p)
            if p not in self.files:
                self.files.append(p)
        self._refresh_list()
        self.path_var.set(f"{len(self.files)} file(s) queued")

    def _refresh_list(self) -> None:
        self.listbox.delete(0, tk.END)
        for p in self.files:
            self.listbox.insert(tk.END, os.path.basename(p))

    def _lb_press(self, event: tk.Event) -> None:
        if self._busy:
            return
        idx = self.listbox.nearest(event.y)
        if 0 <= idx < len(self.files):
            self._drag_index = idx
            self.listbox.selection_clear(0, tk.END)
            self.listbox.selection_set(idx)

    def _lb_motion(self, event: tk.Event) -> None:
        if self._drag_index is None or self._busy:
            return
        cur = self.listbox.nearest(event.y)
        if cur < 0 or cur >= len(self.files) or cur == self._drag_index:
            return
        item = self.files.pop(self._drag_index)
        self.files.insert(cur, item)
        self._refresh_list()
        self.listbox.selection_set(cur)
        self._drag_index = cur

    def _lb_release(self, event: tk.Event) -> None:
        self._drag_index = None

    def _remove_selected(self) -> None:
        sel = self.listbox.curselection()
        if not sel:
            return
        i = int(sel[0])
        self.files.pop(i)
        self._refresh_list()
        self.path_var.set(f"{len(self.files)} file(s) queued")

    def _clear_all(self) -> None:
        self.files.clear()
        self._refresh_list()
        self.path_var.set("")

    def _parse_optional_float(self, raw: str) -> float | None:
        s = raw.strip()
        if not s:
            return None
        return float(s)

    def _parse_optional_font_size(self, raw: str) -> float | None:
        s = raw.strip()
        if not s:
            return None
        v = float(s)
        if v <= 0:
            raise ValueError("Font size must be a positive number.")
        return v

    def _export(self) -> None:
        if self._busy:
            return
        if not self.files:
            messagebox.showwarning("No files", "Add at least one .pptx file.")
            return
        out = filedialog.asksaveasfilename(
            title="Save fixed presentation",
            defaultextension=".pptx",
            filetypes=[("PowerPoint", "*.pptx")],
        )
        if not out:
            return
        try:
            ls = self._parse_optional_float(self.ent_line_spc.get())
            if ls is not None and ls <= 0:
                raise ValueError("Line spacing must be a positive number.")
            fs = self._parse_optional_font_size(self.ent_font_size.get())
        except ValueError as e:
            messagebox.showerror("Invalid input", str(e))
            return

        persian = self.cb_persian.get().strip() or DEFAULT_PERSIAN
        latin = self.cb_latin.get().strip() or DEFAULT_LATIN

        w_cfg = None
        if self.var_watermark.get():
            w_cfg = WatermarkRemovalConfig(
                "strict",
                None,
                "Image 0",
                (8029180, 4844491, 1071843, 256032),
                8192,
                0.55,
                0.55,
                None,
                False,
            )

        self._set_busy(True)
        paths_snapshot = list(self.files)

        def job() -> None:
            tmp_path: str | None = None
            try:
                if len(paths_snapshot) == 1:
                    src = paths_snapshot[0]
                else:
                    fd, tmp_path = tempfile.mkstemp(suffix=".pptx")
                    os.close(fd)
                    merge_pptx_files(paths_snapshot, tmp_path)
                    src = tmp_path

                transform_pptx(
                    src,
                    out,
                    target_cx=12192000,
                    target_cy=6858000,
                    typography=True,
                    persian_font=persian,
                    latin_font=latin,
                    latin_lang="en-US",
                    persian_lang="fa-IR",
                    cs_on_all_runs=True,
                    set_ea_font=True,
                    font_size=fs,
                    strip_margins=False,
                    line_spacing_single=False,
                    line_spacing_multiple=ls,
                    drop_sldsz_type=True,
                    parts=parse_parts_arg(DEFAULT_PARTS),
                    watermark_cfg=w_cfg,
                )
            except Exception as e:  # noqa: BLE001 — surface to user
                # Bind str(e) now: `e` is deleted after `except`, so lambdas must not close over `e`.
                self.master.after(
                    0,
                    lambda msg=str(e): messagebox.showerror("Export failed", msg),
                )
            else:
                self.master.after(
                    0,
                    lambda: messagebox.showinfo(
                        "Done", f"Saved:\n{out}"
                    ),
                )
            finally:
                if tmp_path and os.path.isfile(tmp_path):
                    try:
                        os.remove(tmp_path)
                    except OSError:
                        pass
                self.master.after(0, lambda: self._set_busy(False))

        threading.Thread(target=job, daemon=True).start()

    def _set_busy(self, busy: bool) -> None:
        self._busy = busy
        self.btn_export.config(state=tk.DISABLED if busy else tk.NORMAL)


def main() -> None:
    root = tk.Tk()
    style = ttk.Style()
    if "vista" in style.theme_names():
        style.theme_use("vista")
    PptxLayoutApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
